from __future__ import annotations

import json
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------- Style ----------

BG = RGBColor(6, 11, 24)  # deep navy
CARD = RGBColor(17, 24, 39)  # slate-ish
PRIMARY = RGBColor(255, 203, 48)  # gold
TEXT = RGBColor(245, 247, 255)  # near-white
TEXT_DIM = RGBColor(180, 190, 210)
BORDER = RGBColor(70, 85, 110)
ACCENT_BLUE = RGBColor(59, 130, 246)

FONT_HEAD = "Segoe UI"
FONT_BODY = "Segoe UI"


@dataclass(frozen=True)
class ScreenshotSpec:
    id: str
    title: str
    image: str


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _load_manifest(manifest_path: Path) -> tuple[str, Path, Path, list[ScreenshotSpec]]:
    data = json.loads(manifest_path.read_text(encoding="utf-8"))
    deck_title = str(data.get("deckTitle", "IPL Auction Platform"))
    output_file = _repo_root() / str(data.get("outputFile", "ppt/output/IPL_Auction_Website_Overview.pptx"))
    screenshots_dir = _repo_root() / str(data.get("screenshotsDir", "ppt/assets/screenshots"))
    slides_raw = data.get("slides", [])

    slides: list[ScreenshotSpec] = []
    for item in slides_raw:
        slides.append(ScreenshotSpec(id=item["id"], title=item["title"], image=item["image"]))

    return deck_title, output_file, screenshots_dir, slides


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title(slide, title: str, subtitle: Optional[str] = None) -> None:
    # Title band
    x, y, w, h = Inches(0.8), Inches(0.6), Inches(12.0), Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    p.alignment = PP_ALIGN.LEFT
    run.font.name = FONT_HEAD
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = PRIMARY

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.level = 0
        p2.alignment = PP_ALIGN.LEFT
        p2.font.name = FONT_BODY
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_DIM


def _add_footer(slide, text: str) -> None:
    x, y, w, h = Inches(0.8), Inches(7.05), Inches(12.0), Inches(0.35)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DIM


def _add_card(slide, x, y, w, h, title: str) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)

    # card title
    tb = slide.shapes.add_textbox(x + Inches(0.35), y + Inches(0.25), w - Inches(0.7), Inches(0.4))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEAD
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY


def _add_bullets(slide, x, y, w, h, bullets: Iterable[str]) -> None:
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT


def _add_screenshot_or_placeholder(
    slide,
    image_path: Path,
    x,
    y,
    w,
    h,
    placeholder_label: str,
) -> None:
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), x, y, w, h)
        return

    # placeholder box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0, 0, 0)
    box.fill.transparency = 0.85
    box.line.color.rgb = ACCENT_BLUE
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Add Screenshot"
    run.font.name = FONT_HEAD
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = TEXT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = placeholder_label
    p2.font.name = FONT_BODY
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_DIM


def _set_notes(slide, notes: str) -> None:
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.clear()
    tf.text = notes


def _new_slide(prs: Presentation) -> object:
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def build_deck(deck_title: str, screenshots_dir: Path, specs: list[ScreenshotSpec], output_file: Path) -> None:
    prs = Presentation()

    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title ----
    s = _new_slide(prs)
    _set_bg(s, BG)
    _add_title(s, deck_title, "Real-time IPL player auction website — feature walkthrough")

    hero = s.shapes.add_textbox(Inches(0.9), Inches(2.25), Inches(12.0), Inches(2.2))
    tf = hero.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Create a room · Invite friends · Bid live · Build squads · Export results"
    p.font.name = FONT_BODY
    p.font.size = Pt(26)
    p.font.color.rgb = TEXT

    p2 = tf.add_paragraph()
    p2.text = "Includes: Manual setup, Paddle mode, Watchlist, Icon picks, Voice room, Spectator view"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_DIM

    _set_notes(
        s,
        "Animation: Title Fade (0.3s).\n"
        "Tip: Use a subtle Zoom on the hero line (optional).",
    )

    # ---- Slide 2: Agenda ----
    s = _new_slide(prs)
    _set_bg(s, BG)
    _add_title(s, "Agenda")
    _add_card(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(4.9), "What we will cover")
    _add_bullets(
        s,
        Inches(1.3),
        Inches(2.55),
        Inches(10.8),
        Inches(4.0),
        [
            "Home flow (Create / Join / Watch Live)",
            "Lobby (invite links, watchlist, icon players)",
            "Auction screen (bidding, chat, host controls)",
            "Paddle mode (host bids for all teams)",
            "SOLD NOW (instant finalize)",
            "Voice room (live audio via WebRTC)",
            "Results + export (PDF)",
        ],
    )
    _set_notes(s, "Animation: Reveal bullets one-by-one (Fade, 0.2s).")

    # ---- Slide 3: Key highlights ----
    s = _new_slide(prs)
    _set_bg(s, BG)
    _add_title(s, "Key Highlights")
    _add_card(s, Inches(0.9), Inches(1.85), Inches(5.8), Inches(5.1), "Why players love it")
    _add_bullets(
        s,
        Inches(1.25),
        Inches(2.5),
        Inches(5.1),
        Inches(4.3),
        [
            "Real-time bidding (instant sync)",
            "Mobile-friendly responsive UI",
            "Private rooms with passcode",
            "Watchlist + Icon players",
            "Spectator mode (watch live)",
            "Voice chat inside the auction",
        ],
    )

    _add_card(s, Inches(7.0), Inches(1.85), Inches(5.5), Inches(5.1), "Host power features")
    _add_bullets(
        s,
        Inches(7.35),
        Inches(2.5),
        Inches(4.8),
        Inches(4.3),
        [
            "Manual auction setup",
            "Host Manager mode (manage-only)",
            "Paddle mode (host bids for everyone)",
            "Pause / terminate controls",
            "SOLD NOW to save time",
            "Export results PDF",
        ],
    )

    _set_notes(s, "Animation: Two cards appear (Fade). Then bullets (Fade).")

    # ---- Screenshot-driven slides ----
    for idx, spec in enumerate(specs, start=1):
        s = _new_slide(prs)
        _set_bg(s, BG)
        _add_title(s, f"{idx:02d}. {spec.title}")

        img = screenshots_dir / spec.image
        _add_screenshot_or_placeholder(
            s,
            img,
            Inches(0.9),
            Inches(1.75),
            Inches(11.6),
            Inches(5.15),
            f"File: {spec.image}\\nFolder: {screenshots_dir.as_posix()}",
        )

        notes_map = {
            "home": "Animation: Morph from Title → Home slide (optional).",
            "auth": "Animation: Fade in screenshot. Then highlight Login/Signup with a rectangle (manual).",
            "create": "Animation: Appear (0.2s). Talk through Budget / Squad / Timer / Passcode.",
            "join": "Animation: Appear. Mention ‘Watch Live’ for spectators.",
            "manual": "Animation: Appear. Emphasize Host Role modes: Playing / Manager / Paddle.",
            "lobby": "Animation: Appear. Show Copy link, WhatsApp share, Start button.",
            "watchlist": "Animation: Appear. Explain star favorites + highlight when player appears.",
            "iconpick": "Animation: Appear. Explain fixed-price pre-picks (icons).",
            "auction": "Animation: Appear. Explain 3-column UI and no-scroll design.",
            "bidding": "Animation: Appear. Explain Base bid, increments, Withdraw, Pass, Purse warning.",
            "paddle": "Animation: Appear. Explain host selects acting team then bids.",
            "soldnow": "Animation: Appear. Explain instant finalize when a highest bid exists.",
            "voice": "Animation: Appear. Explain Join Voice, mute/unmute, host mute.",
            "results": "Animation: Appear. Explain team squads + export PDF.",
        }
        _set_notes(s, notes_map.get(spec.id, "Animation: Fade in screenshot (0.3s)."))

        _add_footer(s, "IPL Auction · 2026")

    # ---- Closing slide ----
    s = _new_slide(prs)
    _set_bg(s, BG)
    _add_title(s, "Thank You")
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(12.0), Inches(3.0))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Ready to host your next auction?"
    p.font.name = FONT_HEAD
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "Open the website → Create room → Share code → Start bidding"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT

    _set_notes(s, "Animation: Fade. Add a final click sound (optional).")

    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))


def main() -> None:
    root = _repo_root()
    manifest_path = root / "ppt/assets/manifest.json"
    if not manifest_path.exists():
        raise SystemExit(f"Missing manifest: {manifest_path}")

    deck_title, output_file, screenshots_dir, specs = _load_manifest(manifest_path)
    screenshots_dir.mkdir(parents=True, exist_ok=True)

    build_deck(deck_title, screenshots_dir, specs, output_file)
    print(f"✅ PPT generated: {output_file}")


if __name__ == "__main__":
    main()
